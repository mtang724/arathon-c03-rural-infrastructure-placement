"""
Stage 12 -- the six-slide deck.

Everything is a native PowerPoint object: real chart parts (editable in
PowerPoint, with their own data sheets), autoshapes, connectors and tables.
No images, so every figure stays vector and stays editable.

The scatter "maps" are genuine XY charts with two series, which is what lets a
map be native rather than a picture.
"""
import json

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from config import DATA, REPORTS, ROOT, SERVING_SITE
from coverage import avail_to_rsrp, build_grid, Scorer
from coverage_terrain import fit_with_terrain
from features import haversine_m, load_sites
from model import fit_outage_curve, fit_pathloss

# ---- palette (the validated one used throughout the project) --------------
INK   = RGBColor(0x16, 0x21, 0x1C)
INK2  = RGBColor(0x3A, 0x48, 0x42)
MUTE  = RGBColor(0x65, 0x72, 0x6B)
BG    = RGBColor(0xFB, 0xFC, 0xFA)
SURF  = RGBColor(0xF1, 0xF4, 0xEF)
RULE  = RGBColor(0xC3, 0xCC, 0xBF)
TEAL  = RGBColor(0x0F, 0x6E, 0x70)
OCHRE = RGBColor(0x8F, 0x62, 0x00)
VIOL  = RGBColor(0x5B, 0x3A, 0x9B)
WINE  = RGBColor(0x8C, 0x1D, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
FD, FM = "Archivo", "Consolas"


# ==========================================================================
# helpers
# ==========================================================================

def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=14, color=INK, bold=False, font=FD,
        align=PP_ALIGN.LEFT, caps=False, space=0, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text.upper() if caps else text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = color
    if space:
        # letter-spacing lives on a:rPr as a plain (un-namespaced) attribute
        r.font._rPr.set("spc", str(int(space * 100)))
    return tb


def rect(s, x, y, w, h, fill=SURF, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def oval(s, x, y, w, h, fill=None, line=TEAL, lw=1.25):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def header(s, num, eyebrow, title, sub=None):
    rect(s, Inches(0), Inches(0), W, Inches(0.055), fill=TEAL)
    txt(s, Inches(0.55), Inches(0.34), Inches(0.6), Inches(0.3),
        f"{num:02d}", 12, TEAL, True, FM)
    txt(s, Inches(1.05), Inches(0.36), Inches(6), Inches(0.3),
        eyebrow, 9.5, MUTE, False, FM, caps=True, space=1.6)
    txt(s, Inches(0.55), Inches(0.68), Inches(11.5), Inches(0.62),
        title, 30, INK, True, FD)
    if sub:
        txt(s, Inches(0.55), Inches(1.28), Inches(11.5), Inches(0.32),
            sub, 12.5, INK2, False, FD)


def kpi(s, x, y, w, label, value, note, vcolor=INK, h=Inches(1.02)):
    rect(s, x, y, w, h, fill=SURF, line=RULE)
    txt(s, x + Inches(0.14), y + Inches(0.10), w - Inches(0.28), Inches(0.16),
        label, 8, MUTE, False, FM, caps=True, space=1.2)
    txt(s, x + Inches(0.14), y + Inches(0.28), w - Inches(0.28), Inches(0.42),
        value, 25, vcolor, True, FD)
    txt(s, x + Inches(0.14), y + Inches(0.72), w - Inches(0.28), Inches(0.24),
        note, 8.5, MUTE, False, FM)


def style_chart(ch, size=9, legend=False, gridlines=False):
    ch.font.size = Pt(size); ch.font.name = FM; ch.font.color.rgb = MUTE
    ch.has_title = False
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(size); ch.legend.font.name = FM
        ch.legend.font.color.rgb = INK2
    else:
        ch.has_legend = False
    try:
        va = ch.value_axis
        va.has_major_gridlines = gridlines
        if gridlines:
            gl = va.major_gridlines.format.line
            gl.color.rgb = RULE; gl.width = Pt(0.5)
        va.format.line.color.rgb = RULE
        va.major_tick_mark = XL_TICK_MARK.NONE
        va.tick_labels.font.size = Pt(size - 0.5)
        va.tick_labels.font.name = FM; va.tick_labels.font.color.rgb = MUTE
    except (ValueError, AttributeError):
        pass
    try:
        ca = ch.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = RULE
        ca.major_tick_mark = XL_TICK_MARK.NONE
        ca.tick_labels.font.size = Pt(size - 0.5)
        ca.tick_labels.font.name = FM; ca.tick_labels.font.color.rgb = MUTE
    except (ValueError, AttributeError):
        pass
    return ch


def bar(s, x, y, w, h, cats, series, colors, horizontal=False, gridlines=True,
        labels=False, numfmt='0.0', size=9, gap=60, overlap=-20):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series:
        cd.add_series(nm, vals)
    t = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = s.shapes.add_chart(t, x, y, w, h, cd)
    ch = gf.chart
    for i, pl in enumerate(ch.plots):
        pl.gap_width = gap
        if len(series) > 1:
            pl.overlap = overlap
        pl.has_data_labels = labels
        if labels:
            dl = pl.data_labels
            dl.font.size = Pt(size - 0.5); dl.font.name = FM
            dl.font.color.rgb = INK; dl.number_format = numfmt
            dl.number_format_is_linked = False
    for i, sr in enumerate(ch.series):
        sr.format.fill.solid(); sr.format.fill.fore_color.rgb = colors[i % len(colors)]
        sr.format.line.fill.background()
    return style_chart(ch, size, legend=len(series) > 1, gridlines=gridlines)


def scatter(s, x, y, w, h, series, colors, sizes=None, legend=True, size=9):
    cd = XyChartData()
    for nm, pts in series:
        sd = cd.add_series(nm)
        for px, py in pts:
            sd.add_data_point(px, py)
    gf = s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, w, h, cd)
    ch = gf.chart
    for i, sr in enumerate(ch.series):
        sr.format.line.fill.background()
        m = sr.marker
        m.style = 8  # circle
        m.size = (sizes or [3] * len(series))[i]
        m.format.fill.solid(); m.format.fill.fore_color.rgb = colors[i]
        m.format.line.fill.background()
    return style_chart(ch, size, legend=legend, gridlines=False)


def line_chart(s, x, y, w, h, cats, series, colors, gridlines=True, size=9,
               markers=False):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series:
        cd.add_series(nm, vals)
    t = XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE
    gf = s.shapes.add_chart(t, x, y, w, h, cd)
    ch = gf.chart
    for i, sr in enumerate(ch.series):
        sr.format.line.color.rgb = colors[i % len(colors)]
        sr.format.line.width = Pt(2.25)
        sr.smooth = False
    return style_chart(ch, size, legend=len(series) > 1, gridlines=gridlines)


def caption(s, x, y, w, text, size=9.5, color=MUTE):
    txt(s, x, y, w, Inches(0.5), text, size, color, False, FM)


# ==========================================================================
# data
# ==========================================================================

def gather():
    df = pd.read_csv(DATA / "labeled_terrain.csv", dtype={"cellid": str})
    plf = fit_pathloss(df); oc = fit_outage_curve(df, plf); pl = fit_with_terrain(df)
    thr = avail_to_rsrp(oc, 0.50)
    cov = json.load((REPORTS / "coverage_terrain.json").open())
    ex = json.load((REPORTS / "deck_extras.json").open())
    rob = json.load((REPORTS / "robustness.json").open())
    df["outage"] = df.cellid.isna() | df.cellid.eq("FFFFFFFFF")
    df["obstr"] = df.fresnel_frac < 0.55

    rings = [(0, 2000), (2000, 4000), (4000, 6000), (6000, 8000), (8000, 13000)]
    ring_lbl, out_pct, obs_clear, obs_blocked, ul_med, dl_med = [], [], [], [], [], []
    for lo, hi in rings:
        m = df[(df.dist_m >= lo) & (df.dist_m < hi)]
        ring_lbl.append(f"{lo//1000}-{hi//1000} km")
        out_pct.append(100 * float(m.outage.mean()))
        obs_clear.append(100 * float(m[~m.obstr].outage.mean()))
        obs_blocked.append(100 * float(m[m.obstr].outage.mean()))
        ul_med.append(float(m.uplink.median()) if m.uplink.notna().any() else 0)
        dl_med.append(float(m.downlink.median()) if m.downlink.notna().any() else 0)

    b0, sl = pl["b0"], pl["slope"]
    radii = {d: 10 ** ((b0 - d - thr) / (-sl)) for d in (0, 20, 26)}

    sites, _ = load_sites(); tl, to = sites[SERVING_SITE]
    sub = df.iloc[::4]
    return dict(df=df, pl=pl, thr=thr, cov=cov, rob=rob, ex=ex, radii=radii,
                ring_lbl=ring_lbl, out_pct=out_pct, obs_clear=obs_clear,
                obs_blocked=obs_blocked, ul=ul_med, dl=dl_med,
                served=[(float(r.lon), float(r.lat)) for r in sub.itertuples() if not r.outage],
                dead=[(float(r.lon), float(r.lat)) for r in sub.itertuples() if r.outage],
                tower=(to, tl))


# ==========================================================================
# slides
# ==========================================================================

def s1(prs, D):
    s = slide(prs)
    header(s, 1, "The problem", "Six in ten kilometres have no usable service",
           "7,144 measurements, two days, 178 km² of Iowa farmland — served by one tower")
    x = Inches(0.55)
    for lab, val, note, col in [
            ("Samples", "7,144", "one every 2.6 s", INK),
            ("No serving cell", "42.3%", "3,023 rows", WINE),
            ("Route driven", "117 km", "distinct road", INK),
            ("Survey box", "178 km²", "11.0 × 16.2 km", INK)]:
        kpi(s, x, Inches(1.72), Inches(1.72), lab, val, note, col)
        x += Inches(1.83)

    txt(s, Inches(0.55), Inches(3.05), Inches(6.6), Inches(0.22),
        "Every measurement, coloured by whether it had service", 10, INK2, True, FD)
    ch = scatter(s, Inches(0.4), Inches(3.3), Inches(7.0), Inches(3.7),
                 [("had service", D["served"]), ("no serving cell", D["dead"])],
                 [TEAL, WINE], sizes=[3, 3])
    ch.value_axis.minimum_scale = 41.93; ch.value_axis.maximum_scale = 42.04
    ch.category_axis.minimum_scale = -93.88; ch.category_axis.maximum_scale = -93.66

    txt(s, Inches(7.7), Inches(3.05), Inches(5.2), Inches(0.22),
        "Service collapses with distance from the tower", 10, INK2, True, FD)
    bar(s, Inches(7.55), Inches(3.3), Inches(5.3), Inches(1.72),
        D["ring_lbl"], [("% of samples with no cell", D["out_pct"])], [WINE],
        labels=True, numfmt='0"%"')
    txt(s, Inches(7.7), Inches(5.15), Inches(5.2), Inches(0.22),
        "Downlink holds. Uplink does not.", 10, INK2, True, FD)
    bar(s, Inches(7.55), Inches(5.38), Inches(5.3), Inches(1.62),
        D["ring_lbl"], [("downlink Mbps", D["dl"]), ("uplink Mbps", D["ul"])],
        [TEAL, OCHRE])
    rect(s, Inches(0.55), Inches(6.62), Inches(12.23), Inches(0.62), fill=SURF, line=TEAL, lw=1.5)
    txt(s, Inches(0.78), Inches(6.76), Inches(11.8), Inches(0.34),
        "Where does one more installation go?   The tower sits 1.6 km from the north edge "
        "and 9.5 km from the south — the survey reaches six times further south than north.",
        12.5, INK, True, FD)


def s2(prs, D):
    s = slide(prs)
    header(s, 2, "Why it is hard · 1 of 2",
           "Between 2 and 6 km, dead spots are terrain",
           "Obstruction of the first Fresnel zone — not blocked line of sight")

    # Fresnel geometry, drawn from native shapes
    gx, gy, gw, gh = Inches(0.55), Inches(1.75), Inches(5.9), Inches(2.5)
    rect(s, gx, gy, gw, gh, fill=SURF, line=RULE)
    base = gy + Inches(1.95)
    rect(s, gx + Inches(0.3), base, Inches(5.3), Inches(0.16), fill=RULE)
    mast = rect(s, gx + Inches(0.55), gy + Inches(0.5), Inches(0.07), Inches(1.45), fill=INK)
    oval(s, gx + Inches(0.44), gy + Inches(0.34), Inches(0.3), Inches(0.3), fill=INK, line=None)
    hill = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                              gx + Inches(2.6), base - Inches(0.62),
                              Inches(1.5), Inches(0.62))
    hill.fill.solid(); hill.fill.fore_color.rgb = RULE; hill.line.fill.background()
    hill.shadow.inherit = False
    oval(s, gx + Inches(4.75), base - Inches(0.34), Inches(0.18), Inches(0.18),
         fill=TEAL, line=None)
    ell = oval(s, gx + Inches(0.58), gy + Inches(0.72), Inches(4.3), Inches(1.1),
               fill=None, line=OCHRE, lw=1.5)
    ln = s.shapes.add_connector(1, gx + Inches(0.6), gy + Inches(0.5),
                               gx + Inches(4.84), base - Inches(0.25))
    ln.line.color.rgb = INK; ln.line.width = Pt(1.5)
    txt(s, gx + Inches(0.25), gy + Inches(0.12), Inches(1.6), Inches(0.2),
        "120 ft mast", 8.5, INK2, True, FM)
    txt(s, gx + Inches(1.9), gy + Inches(0.55), Inches(2.2), Inches(0.2),
        "first Fresnel zone", 8.5, OCHRE, True, FM)
    txt(s, gx + Inches(2.45), base + Inches(0.22), Inches(2.0), Inches(0.2),
        "ridge clips the zone", 8.5, WINE, True, FM)
    txt(s, gx + Inches(4.25), base - Inches(0.62), Inches(1.4), Inches(0.2),
        "receiver", 8.5, TEAL, True, FM)
    txt(s, gx + Inches(0.3), gy + Inches(2.18), Inches(5.4), Inches(0.24),
        "Line of sight is clear. The link is still obstructed.", 10.5, INK, True, FD)

    txt(s, Inches(6.85), Inches(1.75), Inches(6), Inches(0.22),
        "Outage rate, split by whether the path is obstructed", 10, INK2, True, FD)
    bar(s, Inches(6.7), Inches(1.98), Inches(6.1), Inches(2.3), D["ring_lbl"],
        [("path clear", D["obs_clear"]), ("Fresnel obstructed", D["obs_blocked"])],
        [TEAL, WINE], labels=True, numfmt='0"%"')

    x = Inches(0.55)
    for lab, val, note, col in [
            ("Blocked line of sight", "13%", "of all links", MUTE),
            ("Inside the Fresnel zone", "46%", "the real test", OCHRE),
            ("Odds of outage, 2-4 km", "2.25×", "p = 7.7e-08", WINE),
            ("Odds of outage, 4-6 km", "2.39×", "p = 5.5e-12", WINE),
            ("Beyond 6 km", "no effect", "link budget is gone", MUTE)]:
        kpi(s, x, Inches(4.6), Inches(2.34), lab, val, note, col)
        x += Inches(2.45)
    caption(s, Inches(0.55), Inches(5.95), Inches(12.2),
            "Obstructed cells are twice as likely to be dead as clear cells the same distance out. "
            "Past 6 km the effect vanishes — everything fails there anyway.")
    caption(s, Inches(0.55), Inches(6.45), Inches(12.2),
            "Terrain relief across the survey: 98 m. DEM: USGS 3DEP 1/3 arc-second, 10 m posts.", 9.5, INK2)


def s3(prs, D):
    s = slide(prs)
    pl = D["pl"]; bt = D["ex"]["backtest"]
    header(s, 3, "How we model it — and what we checked",
           "One physical chain, and an honest account of where it holds",
           "Geometry to received power to service, then tested against the measurements it did not see")

    stages = [("GEOMETRY", "distance, bearing,\nterrain profile"),
              ("RSRP", "path loss +\nknife-edge diffraction"),
              ("AVAILABILITY", "fitted curve,\nP(service) vs RSRP"),
              ("COVERAGE", "route-km 70%\narea 30%")]
    x = Inches(0.55)
    for i, (t, b) in enumerate(stages):
        col = [TEAL, TEAL, OCHRE, INK][i]
        rect(s, x, Inches(1.68), Inches(2.62), Inches(1.12), fill=SURF, line=col, lw=1.5)
        txt(s, x + Inches(0.16), Inches(1.79), Inches(2.3), Inches(0.2),
            t, 9.5, col, True, FM, caps=True, space=1.2)
        txt(s, x + Inches(0.16), Inches(2.02), Inches(2.3), Inches(0.7),
            b, 10, INK2, False, FD)
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.68),
                                    Inches(2.11), Inches(0.34), Inches(0.26))
            ar.fill.solid(); ar.fill.fore_color.rgb = RULE
            ar.line.fill.background(); ar.shadow.inherit = False
        x += Inches(3.09)
    txt(s, Inches(0.55), Inches(2.86), Inches(6.2), Inches(0.2),
        "The first two links are strong. The third is not.", 11, OCHRE, True, FD)

    txt(s, Inches(0.55), Inches(3.24), Inches(3.8), Inches(0.2),
        "Terrain absorbs a fifth of the scatter", 9.5, INK2, True, FD)
    bar(s, Inches(0.4), Inches(3.44), Inches(3.9), Inches(2.0),
        ["distance\n+ azimuth", "+ terrain\ndiffraction"],
        [("residual σ, dB", [9.21, pl["sigma"]])], [OCHRE], labels=True, numfmt='0.00')

    txt(s, Inches(4.5), Inches(3.24), Inches(3.8), Inches(0.2),
        "and the exponent falls to free space", 9.5, INK2, True, FD)
    bar(s, Inches(4.35), Inches(3.44), Inches(3.9), Inches(2.0),
        ["without\nterrain", "with\nterrain", "free\nspace"],
        [("path-loss exponent", [2.40, pl["n_exponent"], 2.00])],
        [TEAL, TEAL, RULE], labels=True, numfmt='0.00')

    txt(s, Inches(8.45), Inches(3.24), Inches(4.4), Inches(0.2),
        "RSRP error, three ways of splitting the data", 9.5, INK2, True, FD)
    bar(s, Inches(8.3), Inches(3.44), Inches(4.5), Inches(2.0),
        ["in\nsample", "random\nsplit", "held-out\ngeography"],
        [("MAE, dB", [bt["rsrp_mae_in"], bt["rsrp_mae_random"], bt["rsrp_mae_spatial"]])],
        [RULE, RULE, WINE], labels=True, numfmt='0.00')

    for i, (lab, val, note, col) in enumerate([
            ("Diffraction coefficient", f"{pl['b_diff']:.2f}", "dB per dB · physics says −1", INK),
            ("RSRP held-out R²", f"{bt['rsrp_r2_spatial']:+.2f}", "vs +0.79 on a random split", WINE),
            ("Availability, our model", f"{bt['acc_model']:.0f}%", f"base rate is {bt['base_rate']:.0f}%", WINE),
            ("Where vs when", "71 / 5%", "of outage variance", INK)]):
        kpi(s, Inches(0.55) + Inches(3.13) * i, Inches(5.6), Inches(3.0), lab, val, note, col,
            h=Inches(0.92))

    rect(s, Inches(0.55), Inches(6.62), Inches(12.23), Inches(0.66), fill=SURF, line=WINE, lw=1.5)
    txt(s, Inches(0.78), Inches(6.75), Inches(11.8), Inches(0.42),
        "What the backtest says: the propagation model is sound, the availability step is weak. "
        "Received power we predict to 5.9 dB; whether a cell has service we barely beat guessing. "
        "So the device comparison holds and the absolute coverage percentages do not.",
        11.5, INK, True, FD)


def s4(prs, D):
    s = slide(prs)
    r = D["radii"]; cov = D["cov"]["assets"]
    header(s, 4, "Why it is hard · 2 of 2", "No device on the menu can reach",
           "Every 6 dB of transmit power lost halves the radius and quarters the area")

    # concentric circles, drawn to true scale
    cxp, cyp = Inches(3.15), Inches(4.25)
    scale = Inches(2.25).inches / r[0]
    for d, col, nm in [(0, TEAL, "Macro"), (20, OCHRE, "Relay"), (26, WINE, "Small cell")]:
        rad = Inches(r[d] * scale)
        oval(s, cxp - rad, cyp - rad, rad * 2, rad * 2, fill=None, line=col, lw=2.0)
    txt(s, Inches(0.6), Inches(1.72), Inches(5.2), Inches(0.22),
        "Service radius, drawn to scale", 10, INK2, True, FD)
    oval(s, cxp - Inches(0.05), cyp - Inches(0.05), Inches(0.1), Inches(0.1),
         fill=INK, line=None)
    txt(s, cxp + Inches(0.12), cyp - Inches(2.28), Inches(2.4), Inches(0.2),
        f"Macro  {r[0]:,.0f} m", 10, TEAL, True, FM)
    txt(s, cxp + Inches(0.12), cyp - Inches(0.34), Inches(2.4), Inches(0.2),
        f"Relay  {r[20]:.0f} m", 10, OCHRE, True, FM)
    txt(s, cxp + Inches(0.12), cyp - Inches(0.13), Inches(2.4), Inches(0.2),
        f"Small cell  {r[26]:.0f} m", 10, WINE, True, FM)

    rows = [("", "Macro-class", "Donor relay", "Small cell"),
            ("Power vs the tower", "same", "−20 dB", "−26 dB"),
            ("Mast height", "37 m", "10 m", "10 m"),
            ("Backhaul", "fibre + tower", "over the air", "fibre"),
            ("Service radius", f"{r[0]:,.0f} m", f"{r[20]:.0f} m", f"{r[26]:.0f} m"),
            ("Coverage area", "1.00×", "0.009×", "0.002×")]
    tx, ty, tw = Inches(6.35), Inches(1.95), Inches(6.45)
    tbl = s.shapes.add_table(len(rows), 4, tx, ty, tw, Inches(2.35)).table
    tbl.columns[0].width = Inches(2.1)
    for i in range(1, 4):
        tbl.columns[i].width = Inches(1.45)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(0.39)
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = ""
            c.fill.solid(); c.fill.fore_color.rgb = SURF if i == 0 else BG
            c.margin_left = c.margin_right = Inches(0.09)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = v
            f = run.font; f.size = Pt(10.5); f.name = FM if j else FD
            f.bold = (i == 0) or (j == 0)
            f.color.rgb = [INK, TEAL, OCHRE, WINE][j] if i == 0 else (INK2 if j == 0 else INK)

    txt(s, Inches(6.5), Inches(4.5), Inches(3.1), Inches(0.22),
        "Gain collapses with transmit power", 10, INK2, True, FD)
    line_chart(s, Inches(6.2), Inches(4.7), Inches(3.35), Inches(1.98),
               [f"−{d['db']}" for d in D["ex"]["sweep_power"]],
               [("coverage gain", [d["gain"] for d in D["ex"]["sweep_power"]])],
               [WINE], markers=True)
    txt(s, Inches(9.75), Inches(4.5), Inches(3.1), Inches(0.22),
        "What one of each actually adds", 10, INK2, True, FD)
    bar(s, Inches(9.6), Inches(4.7), Inches(3.22), Inches(1.98),
        ["route-km added", "km² added"],
        [("Macro", [cov["macro"]["one_asset"]["route_km_added"],
                    cov["macro"]["one_asset"]["area_km2_added"]]),
         ("Relay", [cov["relay"]["one_asset"]["route_km_added"],
                    cov["relay"]["one_asset"]["area_km2_added"]]),
         ("Small cell", [cov["smallcell"]["one_asset"]["route_km_added"],
                         cov["smallcell"]["one_asset"]["area_km2_added"]])],
        [TEAL, OCHRE, WINE], labels=True, numfmt='0.0')

    caption(s, Inches(0.55), Inches(6.85), Inches(12.2),
            "5 dB below the tower halves the gain. 10 dB cuts it to a sixth. A relay is 20 dB down — a 520 m bubble against a 9 km hole.")


def s5(prs, D):
    s = slide(prs)
    header(s, 5, "Why our solution works · the tool",
           "Everything the brief asks for, in one clickable page",
           "Click anywhere and the terrain model re-solves the whole surface, live in the browser")

    px, py, pw, ph = Inches(0.55), Inches(1.78), Inches(6.9), Inches(4.05)
    rect(s, px, py, pw, ph, fill=SURF, line=RULE)
    for i in range(8):
        ln = s.shapes.add_connector(1, px + Inches(0.28), py + Inches(0.38 + i * 0.44),
                                    px + Inches(6.62), py + Inches(0.58 + i * 0.44))
        ln.line.color.rgb = RULE; ln.line.width = Pt(0.75)
    rng = np.random.default_rng(5)
    for _ in range(80):
        a, b = rng.uniform(0.35, 6.4), rng.uniform(0.28, 3.6)
        oval(s, px + Inches(a), py + Inches(b), Inches(0.05), Inches(0.05),
             fill=(WINE if b > 1.75 else TEAL), line=None)
    oval(s, px + Inches(2.5), py + Inches(1.35), Inches(1.8), Inches(1.8),
         fill=None, line=VIOL, lw=2.0)
    oval(s, px + Inches(3.3), py + Inches(2.15), Inches(0.2), Inches(0.2), fill=VIOL, line=None)
    txt(s, px + Inches(3.55), py + Inches(2.11), Inches(2.4), Inches(0.2),
        "YOUR SITE", 9.5, VIOL, True, FM)
    oval(s, px + Inches(1.3), py + Inches(0.45), Inches(0.16), Inches(0.16), fill=INK, line=None)
    txt(s, px + Inches(1.52), py + Inches(0.42), Inches(2.4), Inches(0.2),
        "existing tower", 9.5, INK, True, FM)
    txt(s, px + Inches(0.24), py + Inches(3.72), Inches(5), Inches(0.2),
        "click to place  ·  drag to pan  ·  scroll to zoom", 9, MUTE, False, FM)

    txt(s, Inches(7.75), Inches(1.78), Inches(5.1), Inches(0.22),
        "The brief asks for four things. Each is a tab.", 10.5, INK, True, FD)
    reqs = [("Thresholds", "before / after at four explicit\nservice definitions", TEAL),
            ("Robustness", "150 fresh shadow-fading draws,\nrun live on the placed site", VIOL),
            ("Gains", "per device class, and per\nsuccessive installation", OCHRE),
            ("Sensitivity", "gain swept against mast height\nand transmit power", WINE)]
    for i, (t, b, col) in enumerate(reqs):
        y = Inches(2.08) + Inches(0.93) * i
        rect(s, Inches(7.6), y, Inches(5.18), Inches(0.82), fill=SURF, line=col, lw=1.5)
        txt(s, Inches(7.78), y + Inches(0.11), Inches(2.0), Inches(0.22),
            t, 11.5, col, True, FD)
        txt(s, Inches(7.78), y + Inches(0.36), Inches(4.8), Inches(0.42),
            b, 9.5, INK2, False, FM)

    txt(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.26),
        "Nothing is a stored answer — the page carries the terrain model and runs the same "
        "physics as the offline optimiser, matching it to about 1%.", 11, INK, True, FD)
    for i, (lab, val, note, col) in enumerate([
            ("Terrain posts carried", "31 m", "3DEP, 501 × 851 grid", INK),
            ("Demand cells re-solved", "4,731", "on every click", INK),
            ("Agreement with the optimiser", "~1%", "verified against the offline run", TEAL),
            ("External dependencies", "none", "one file, opens anywhere", TEAL)]):
        kpi(s, Inches(0.55) + Inches(3.13) * i, Inches(6.38), Inches(3.0), lab, val, note, col,
            h=Inches(0.86))


def s6(prs, D):
    s = slide(prs)
    cov = D["cov"]; a = cov["assets"]; b = cov["baseline"]; ex = D["ex"]
    header(s, 6, "The answer", "One macro-class site, 6.5 km south-west",
           "What one installation of each class can and cannot do")

    txt(s, Inches(0.55), Inches(1.72), Inches(6), Inches(0.22),
        "Route coverage after one installation", 10, INK2, True, FD)
    bar(s, Inches(0.4), Inches(1.94), Inches(6.3), Inches(2.15),
        ["now", "+ small cell", "+ relay", "+ macro"],
        [("% of route-km served",
          [b["route_pct"], a["smallcell"]["one_asset"]["route_pct"],
           a["relay"]["one_asset"]["route_pct"], a["macro"]["one_asset"]["route_pct"]])],
        [RULE, WINE, OCHRE, TEAL], labels=True, numfmt='0.0"%"')

    txt(s, Inches(6.95), Inches(1.72), Inches(6), Inches(0.22),
        "Under four explicit service thresholds", 10, INK2, True, FD)
    rows = [("available", "route now", "after", "area now", "after")] + [
        (f"≥{r['pct']}%", f"{r['route_before']:.1f}%", f"{r['route_after']:.1f}%",
         f"{r['area_before']:.1f}%", f"{r['area_after']:.1f}%") for r in ex["thresholds"]]
    tbl = s.shapes.add_table(len(rows), 5, Inches(6.85), Inches(1.94),
                             Inches(5.95), Inches(2.05)).table
    tbl.columns[0].width = Inches(1.35)
    for i in range(1, 5):
        tbl.columns[i].width = Inches(1.15)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(0.41)
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = ""
            c.fill.solid(); c.fill.fore_color.rgb = SURF if i == 0 else BG
            c.margin_left = c.margin_right = Inches(0.08)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = c.text_frame.paragraphs[0]
            pp.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r_ = pp.add_run(); r_.text = v
            f = r_.font; f.size = Pt(10.5); f.name = FM
            f.bold = (i == 0) or (j == 0)
            hi = (i == 2)   # the 50% row is the headline definition
            f.color.rgb = MUTE if i == 0 else (TEAL if hi and j in (2, 4) else
                                               (INK if j else INK2))

    rect(s, Inches(0.55), Inches(4.25), Inches(6.15), Inches(1.66), fill=SURF, line=TEAL, lw=1.5)
    txt(s, Inches(0.78), Inches(4.42), Inches(5.7), Inches(0.2),
        "The one that works", 9, TEAL, True, FM, caps=True, space=1.4)
    txt(s, Inches(0.78), Inches(4.66), Inches(5.7), Inches(0.34),
        "Macro-class site, 6.5 km south-west", 17, INK, True, FD)
    txt(s, Inches(0.78), Inches(5.06), Inches(5.7), Inches(0.24),
        "41.97956, −93.82983  ·  37 m mast  ·  on the road network", 10.5, INK2, False, FM)
    txt(s, Inches(0.78), Inches(5.36), Inches(5.7), Inches(0.36),
        "40% → 65% of route  ·  34% → 56% of area  ·  1,079 dead cells fixed",
        11.5, TEAL, True, FD)

    txt(s, Inches(6.95), Inches(4.25), Inches(6), Inches(0.22),
        f"Gain over {ex['robust_gain']['n']} shadow-fading draws", 10, INK2, True, FD)
    rg = ex["robust_gain"]
    bar(s, Inches(6.8), Inches(4.45), Inches(3.05), Inches(1.5),
        ["p10", "median", "p90"],
        [("coverage gain", [rg["p10"], rg["p50"], rg["p90"]])],
        [VIOL], labels=True, numfmt='0.000')
    kpi(s, Inches(10.0), Inches(4.47), Inches(2.8), "Draws with a positive gain",
        f"{rg['pct_positive']:.0f}%", "the site is never a mistake", VIOL, h=Inches(1.42))

    for i, (lab, val, note, col) in enumerate([
            ("Relay", "+1.4 pts", "520 m radius", WINE),
            ("Small cell", "+0.9 pts", "256 m radius", WINE),
            ("Still uncovered", "35%", "of route — needs a programme", OCHRE)]):
        kpi(s, Inches(0.55) + Inches(2.09) * i, Inches(6.05), Inches(1.98),
            lab, val, note, col, h=Inches(0.86))
    rect(s, Inches(6.95), Inches(6.05), Inches(5.85), Inches(0.86), fill=SURF, line=WINE, lw=1.5)
    txt(s, Inches(7.15), Inches(6.17), Inches(5.5), Inches(0.64),
        "Read the ratios, not the absolutes. The device comparison rests on the fitted "
        "path-loss law; the percentages inherit a weak availability step that the backtest "
        "showed barely beats the base rate.",
        10.5, INK, True, FD)


def build(verbose=True):
    D = gather()
    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    for fn in (s1, s2, s3, s4, s5, s6):
        fn(prs, D)
    out = ROOT / "ARA_Challenge3.pptx"
    prs.save(out)
    if verbose:
        print(f"[deck] {out} ({out.stat().st_size/1e6:.2f} MB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
