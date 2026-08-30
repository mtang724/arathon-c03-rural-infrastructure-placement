"""
Geometry check for a generated deck.

    python -m common.check_deck ARA_Challenge3.pptx

python-pptx places shapes at absolute coordinates and never reflows, so a text
box that needs more room than it was given does not push anything down -- it
silently draws over whatever is beneath it. That is invisible to any check that
only counts shapes, and it is the failure this looks for.

Three checks:

  OFF-SLIDE   any shape extending past the slide edge.
  OVERFLOW    text that needs more lines than its box has height for. Estimated,
              not measured -- we do not have the font metrics here -- so it is
              deliberately generous and reports the ratio, letting a human judge
              the borderline ones rather than guessing for them.
  COLLISION   two text-bearing shapes whose rectangles intersect. Panels and
              accent bars are meant to sit under their contents, so only
              text-over-text is reported.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
# Archivo is proportional, Consolas is not. Average glyph width as a fraction of
# point size, deliberately on the narrow side so the check errs toward silence.
CHAR_W = {"Consolas": 0.60, "Archivo": 0.50}
LINE_H = 1.35            # line box as a multiple of point size
TOL = 1.18               # report only when a box is over 18% short


def _txt_shapes(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t:
            yield sh, t


def estimate_lines(text, width_in, size_pt, font):
    cw = CHAR_W.get(font, 0.52) * size_pt / 72.0
    per_line = max(1, int(width_in / cw))
    n = 0
    for para in text.split("\n"):
        n += max(1, -(-len(para) // per_line))
    return n


def check(path, verbose=True):
    prs = Presentation(path)
    W = prs.slide_width / EMU_IN
    H = prs.slide_height / EMU_IN
    problems = []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            try:
                x, y = sh.left / EMU_IN, sh.top / EMU_IN
                w, h = sh.width / EMU_IN, sh.height / EMU_IN
            except TypeError:
                continue
            if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
                # the full-bleed background and the accent rule are meant to
                # reach the edges exactly; only flag real overshoot
                if x < -0.01 or y < -0.01 or x + w > W + 0.02 or y + h > H + 0.02:
                    problems.append((idx, "OFF-SLIDE",
                                     f"{sh.shape_type} at ({x:.2f},{y:.2f}) "
                                     f"{w:.2f}x{h:.2f} vs slide {W:.2f}x{H:.2f}"))

        for sh, t in _txt_shapes(slide):
            x, y = sh.left / EMU_IN, sh.top / EMU_IN
            w, h = sh.width / EMU_IN, sh.height / EMU_IN
            runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
            if not runs:
                continue
            size = max((r.font.size.pt for r in runs if r.font.size), default=12)
            font = next((r.font.name for r in runs if r.font.name), "Archivo")
            lines = estimate_lines(t, w, size, font)
            need = lines * size * LINE_H / 72.0
            if need > h * TOL:
                problems.append((idx, "OVERFLOW",
                                 f'"{t[:44]}..." needs ~{need:.2f}in in a '
                                 f'{h:.2f}in box ({need/h:.1f}x) at y={y:.2f}'))
                boxes.append((x, y, w, need, t))          # collide on real need
            else:
                boxes.append((x, y, w, h, t))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ax, ay, aw, ah, at = boxes[i]
                bx, by, bw, bh, bt = boxes[j]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.05 and oy > 0.05:
                    problems.append((idx, "COLLISION",
                                     f'"{at[:26]}..." x "{bt[:26]}..." '
                                     f'overlap {ox:.2f}x{oy:.2f}in at y={max(ay,by):.2f}'))

    if verbose:
        if not problems:
            print(f"[deck] {path}: {len(prs.slides)} slides, no geometry problems")
        else:
            by_slide = {}
            for s, kind, msg in problems:
                by_slide.setdefault(s, []).append((kind, msg))
            for s in sorted(by_slide):
                print(f"\nslide {s}")
                for kind, msg in by_slide[s]:
                    print(f"  {kind:<10} {msg}")
            print(f"\n{len(problems)} problem(s) across "
                  f"{len(by_slide)} of {len(prs.slides)} slides")
    return problems


if __name__ == "__main__":
    p = sys.argv[1] if len(sys.argv) > 1 else "ARA_Challenge3.pptx"
    sys.exit(1 if check(Path(p)) else 0)
