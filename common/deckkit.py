"""
Native-PowerPoint drawing kit.

Everything a deck built with this is made of is a real PowerPoint object -- chart
parts with their own data sheets, autoshapes, connectors, tables. No images, so
every figure stays vector, stays editable, and survives being re-themed by
whoever presents it.

Lifted from terrain-approach's six-slide builder and made model-agnostic so the
repository-wide deck can use it. Nothing here imports an approach.
"""
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- palette -------------------------------------------------------------
INK = RGBColor(0x16, 0x21, 0x1C)
INK2 = RGBColor(0x3A, 0x48, 0x42)
MUTE = RGBColor(0x65, 0x72, 0x6B)
BG = RGBColor(0xFB, 0xFC, 0xFA)
SURF = RGBColor(0xF1, 0xF4, 0xEF)
RULE = RGBColor(0xC3, 0xCC, 0xBF)
TEAL = RGBColor(0x0F, 0x6E, 0x70)
OCHRE = RGBColor(0x8F, 0x62, 0x00)
VIOL = RGBColor(0x5B, 0x3A, 0x9B)
WINE = RGBColor(0x8C, 0x1D, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xA8, 0xB2, 0xA6)

W, H = Inches(13.333), Inches(7.5)
FD, FM = "Archivo", "Consolas"


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=14, color=INK, bold=False, font=FD,
        align=PP_ALIGN.LEFT, caps=False, space=0, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text.upper() if caps else text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = color
    if space:
        r.font._rPr.set("spc", str(int(space * 100)))
    return tb


def bullets(s, x, y, w, h, items, size=12.5, color=INK2, font=FD, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        bold = False
        if isinstance(it, tuple):
            it, bold = it
        r = p.add_run(); r.text = "— " + it
        f = r.font; f.size = Pt(size); f.name = font
        f.color.rgb = INK if bold else color; f.bold = bold
    return tb


def rect(s, x, y, w, h, fill=SURF, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def oval(s, x, y, w, h, fill=None, line=TEAL, lw=1.25):
    return rect(s, x, y, w, h, fill, line, lw, MSO_SHAPE.OVAL)


def arrow(s, x, y, w, h, fill=RULE):
    return rect(s, x, y, w, h, fill, None, 1.0, MSO_SHAPE.RIGHT_ARROW)


def header(s, num, eyebrow, title, sub=None, accent=TEAL):
    rect(s, Inches(0), Inches(0), W, Inches(0.055), fill=accent)
    txt(s, Inches(0.55), Inches(0.34), Inches(0.6), Inches(0.3),
        f"{num:02d}", 12, accent, True, FM)
    txt(s, Inches(1.05), Inches(0.36), Inches(8), Inches(0.3),
        eyebrow, 9.5, MUTE, False, FM, caps=True, space=1.6)
    txt(s, Inches(0.55), Inches(0.68), Inches(12.2), Inches(0.62),
        title, 29, INK, True, FD)
    if sub:
        txt(s, Inches(0.55), Inches(1.26), Inches(12.2), Inches(0.4),
            sub, 12.5, INK2, False, FD)


def kpi(s, x, y, w, label, value, note, vcolor=INK, h=Inches(1.02)):
    rect(s, x, y, w, h, fill=SURF, line=RULE)
    txt(s, x + Inches(0.14), y + Inches(0.10), w - Inches(0.28), Inches(0.16),
        label, 8, MUTE, False, FM, caps=True, space=1.2)
    txt(s, x + Inches(0.14), y + Inches(0.28), w - Inches(0.28), Inches(0.42),
        value, 25, vcolor, True, FD)
    txt(s, x + Inches(0.14), y + Inches(0.72), w - Inches(0.28), Inches(0.24),
        note, 8.5, MUTE, False, FM)


def card(s, x, y, w, h, title, lines, accent=TEAL, status=None, dim=False):
    """A panel for one simulator, or one idea. `dim` marks it as not yet built."""
    rect(s, x, y, w, h, fill=SURF if not dim else BG,
         line=RULE if dim else accent, lw=1.0 if dim else 1.5)
    rect(s, x, y, Inches(0.06), h, fill=GREY if dim else accent)
    txt(s, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.28),
        title, 13.5, MUTE if dim else INK, True, FD)
    if status:
        txt(s, x + Inches(0.22), y + Inches(0.42), w - Inches(0.4), Inches(0.2),
            status, 8.5, GREY if dim else accent, True, FM, caps=True, space=1.1)
    bullets(s, x + Inches(0.22), y + Inches(0.70), w - Inches(0.44),
            h - Inches(0.84), lines, size=10, color=MUTE if dim else INK2, gap=3)
    return None


def style_chart(ch, size=9, legend=False, gridlines=False):
    ch.font.size = Pt(size); ch.font.name = FM; ch.font.color.rgb = MUTE
    ch.has_title = False
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(size); ch.legend.font.name = FM
        ch.legend.font.color.rgb = INK2
    else:
        ch.has_legend = False
    for axis, grid in ((getattr(ch, "value_axis", None), gridlines),
                       (getattr(ch, "category_axis", None), False)):
        try:
            axis.has_major_gridlines = grid
            if grid:
                gl = axis.major_gridlines.format.line
                gl.color.rgb = RULE; gl.width = Pt(0.5)
            axis.format.line.color.rgb = RULE
            axis.major_tick_mark = XL_TICK_MARK.NONE
            axis.tick_labels.font.size = Pt(size - 0.5)
            axis.tick_labels.font.name = FM
            axis.tick_labels.font.color.rgb = MUTE
        except (ValueError, AttributeError):
            pass
    return ch


def bar(s, x, y, w, h, cats, series, colors, horizontal=False, gridlines=True,
        labels=False, numfmt="0.0", size=9, gap=60, overlap=-20):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series:
        cd.add_series(nm, vals)
    t = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    ch = s.shapes.add_chart(t, x, y, w, h, cd).chart
    for pl in ch.plots:
        pl.gap_width = gap
        if len(series) > 1:
            pl.overlap = overlap
        pl.has_data_labels = labels
        if labels:
            dl = pl.data_labels
            dl.font.size = Pt(size - 0.5); dl.font.name = FM
            dl.font.color.rgb = INK; dl.number_format = numfmt
            dl.number_format_is_linked = False
    for i, sr in enumerate(ch.series):
        sr.format.fill.solid(); sr.format.fill.fore_color.rgb = colors[i % len(colors)]
        sr.format.line.fill.background()
    return style_chart(ch, size, legend=len(series) > 1, gridlines=gridlines)


def scatter(s, x, y, w, h, series, colors, sizes=None, legend=True, size=9):
    cd = XyChartData()
    for nm, pts in series:
        sd = cd.add_series(nm)
        for px, py in pts:
            sd.add_data_point(px, py)
    ch = s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, w, h, cd).chart
    for i, sr in enumerate(ch.series):
        sr.format.line.fill.background()
        m = sr.marker
        m.style = 8
        m.size = (sizes or [3] * len(series))[i]
        m.format.fill.solid(); m.format.fill.fore_color.rgb = colors[i]
        m.format.line.fill.background()
    return style_chart(ch, size, legend=legend, gridlines=False)


def line_chart(s, x, y, w, h, cats, series, colors, gridlines=True, size=9,
               markers=False):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series:
        cd.add_series(nm, vals)
    t = XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE
    ch = s.shapes.add_chart(t, x, y, w, h, cd).chart
    for i, sr in enumerate(ch.series):
        sr.format.line.color.rgb = colors[i % len(colors)]
        sr.format.line.width = Pt(2.25)
        sr.smooth = False
    return style_chart(ch, size, legend=len(series) > 1, gridlines=gridlines)


def table(s, x, y, w, h, rows, col_w=None, size=9.5, head_fill=SURF):
    nr, nc = len(rows), len(rows[0])
    tb = s.shapes.add_table(nr, nc, x, y, w, h).table
    if col_w:
        for i, cw in enumerate(col_w):
            tb.columns[i].width = cw
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tb.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = head_fill if r == 0 else WHITE
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            f = run.font
            f.size = Pt(size); f.name = FM; f.bold = (r == 0)
            f.color.rgb = MUTE if r == 0 else INK
            if c > 0:
                p.alignment = PP_ALIGN.RIGHT
    return tb


def caption(s, x, y, w, text, size=9.5, color=MUTE):
    return txt(s, x, y, w, Inches(0.5), text, size, color, False, FM)


def footer(s, text):
    txt(s, Inches(0.55), H - Inches(0.42), Inches(12.2), Inches(0.25),
        text, 8.5, GREY, False, FM)
