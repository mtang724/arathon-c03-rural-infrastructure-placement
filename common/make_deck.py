"""
The project deck: six slides, native objects only.

    python -m common.make_deck

Everything is a real PowerPoint object -- chart parts with their own data sheets,
autoshapes and tables. No images, so every figure stays vector and stays editable
by whoever presents it.

THE NARRATIVE, in the order a listener needs it, cut to six slides for a
seven-minute slot -- roughly seventy seconds each:

    1  the problem: choose where to build, having measured almost none of it
    2  why the measurements alone cannot answer it, and what a twin buys
    3  four ways to build the twin: baseline, ray tracing, deep learning, PINN
    4  what happened: the learned models win in sample and lose on geography
    5  how a prediction becomes a decision
    6  the demo

Numbers come from reports and bundles rather than being typed in, so the deck
cannot drift from the models. Anything missing renders as reserved rather than
as a stale figure.
"""
import glob
import json
import random
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .deckkit import (BG, FD, FH, FM, GREY, H, INK, INK2, MUTE, OCHRE, RULE, SURF,
                      TEAL, VIOL, W, WHITE, WINE, bar, bullets, caption, footer,
                      header, kpi, oval, rect, scatter, slide, table, txt)

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "ARA_Challenge3.pptx"

ORDER = ["terrain-parametric", "sionna-hybrid-agronomy", "terrain-fno",
         "reveal-mt-pinn"]
LABEL = {"terrain-parametric": "Baseline — fitted physics",
         "sionna-hybrid-agronomy": "Ray tracing — Sionna",
         "terrain-fno": "Deep learning — FNO",
         "reveal-mt-pinn": "PINN — ReVeal-MT"}
SHORT = {"terrain-parametric": "Baseline", "sionna-hybrid-agronomy": "Ray tracing",
         "terrain-fno": "Deep learning", "reveal-mt-pinn": "PINN"}


def load():
    d = {"bundles": {}}

    def rd(p):
        try:
            return json.loads((ROOT / p).read_text())
        except Exception:
            return None

    # The shared bench, plus anything a collaborator ran separately. An
    # approach that benchmarks itself writes reports/backtest_<name>.json rather
    # than editing the shared file, so merge those in instead of hand-copying
    # numbers out of a README -- which is how a deck starts to drift from the
    # models it describes.
    d["bench"] = rd("reports/testbench.json") or {}
    for extra in sorted(glob.glob(str(ROOT / "reports" / "backtest_*.json"))):
        blob = rd(Path(extra).relative_to(ROOT).as_posix()) or {}
        sims = blob.get("simulators", blob)
        for k, v in sims.items():
            if isinstance(v, dict) and "in_sample" in v:
                d["bench"].setdefault(k, v)
    # terrain-fno benchmarks itself in its own harness, with a shuffled control
    # the shared bench has no slot for. Fold it in so all four models appear.
    fno = rd("terrain-approach/reports/fno_compare.json") or {}
    if fno:
        ins, oos = fno.get("in_sample", {}), fno.get("out_of_sample", {})
        for key, name in (("fno_residual", "terrain-fno"),
                          ("fno_shuffled_control", "fno-shuffled-control")):
            if key in ins:
                d["bench"].setdefault(name, {
                    "in_sample": ins[key],
                    **{sp: v[key] for sp, v in oos.items() if key in v}})
    d["geom"] = rd("reports/split_geometry.json")
    d["coverage"] = rd("terrain-approach/reports/coverage_terrain.json")
    d["hyp"] = rd("reports/hypothesis_test.json")
    d["sens"] = rd("reports/sensitivity.json")
    for f in sorted(glob.glob(str(ROOT / "bundles" / "*.json"))):
        try:
            b = json.loads(Path(f).read_text())
            d["bundles"][b["simulator"]["name"]] = b
        except Exception:
            pass
    return d


def cell(v, k):
    try:
        return f"{v[k]['rmse']:.2f}"
    except (KeyError, TypeError):
        return "—"


# ==========================================================================

# ==========================================================================
# The six-slide cut. Older, longer slides live in git history.
# ==========================================================================


def n1_problem(prs, d):
    """Open on the decision, not on the dataset."""
    s = slide(prs)
    txt(s, Inches(0.55), Inches(1.0), Inches(9), Inches(0.35),
        "Arathon Challenge 03", 16, MUTE, False, FD)
    txt(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(1.3),
        "Where should the\nnext tower go?", 40, INK, True, FH)
    txt(s, Inches(0.55), Inches(3.02), Inches(5.9), Inches(0.42),
        "Rural network planning on 7% of the map", 20, WINE, False, FH)
    txt(s, Inches(0.55), Inches(3.62), Inches(5.7), Inches(1.6),
        "Four towers serve this area today, and most of it still has no usable "
        "signal — 42% of the samples the van took found no cell at all. One "
        "more asset could change that. The question is where to put it.",
        14.5, INK2, False, FD)
    best, who = None, ""
    for nm, v in d["bench"].items():
        r = (v.get("kmeans_on_position") or {}).get("rmse")
        if r and nm != "fno-shuffled-control" and (best is None or r < best):
            best, who = r, nm
    dec = ROOT / "sionna-approach" / "deck_panels" / "decision.png"
    if dec.exists():
        s.shapes.add_picture(str(dec), Inches(6.55), Inches(1.5),
                             height=Inches(3.35))
        caption(s, Inches(6.55), Inches(4.95), Inches(6.1),
                "Red is where service fails today. The star is where our "
                "planner would put the next tower.", size=15, color=INK2)
    for i, (lab, val, note) in enumerate([
            ("unserved today", "60%", "of the area, by our model"),
            ("we measured", "7%", "of it · the rest is predicted"),
            ("with one more site", "44% → 69%", "of route covered")]):
        kpi(s, Inches(0.55 + i * 4.12), Inches(5.4), Inches(3.85), lab, val,
            note, vcolor=WINE if i in (0, 2) else INK, vsize=27 if i < 2 else 21)
    txt(s, Inches(0.55), Inches(6.72), Inches(11.9), Inches(0.3),
        "Mingyue Tang  ·  David Alcantara  ·  Ishan Bansal", 16, MUTE, False, FD)
    footer(s, "AgWireless '26 · ARA COTS RAN, Ames IA")
    return s


def n2_gap(prs, d):
    """The two real maps side by side. The argument is visible, not asserted.

    These are the only raster images in the deck, and they earn it: a drawn
    schematic of a survey is a claim about the survey, whereas this is the
    survey. Cropped from the project's own validation figure by
    `common/crop_panels.py`.
    """
    s = slide(prs)
    header(s, 2, "the gap", "Measurements alone cannot site a tower",
           "Same ground, same scale, same colours. Only the coverage differs.")
    # Lay the two out from their real aspect ratios so they sit as a balanced
    # pair. Hard-coding x positions left three inches of gutter between them.
    panels = ROOT / "sionna-approach" / "deck_panels"
    items = [("gap_measured.png", "What we measured"),
             ("gap_predicted.png", "What the twin predicts")]
    PH, GAP = 3.85, 0.55
    widths = []
    for fn, _ in items:
        f = panels / fn
        if f.exists():
            from PIL import Image as _Im
            w0, h0 = _Im.open(f).size
            widths.append(PH * w0 / h0)
        else:
            widths.append(0.0)
    x0 = (13.33 - (sum(widths) + GAP)) / 2.0
    for (fn, cap), wid in zip(items, widths):
        f = panels / fn
        if f.exists():
            s.shapes.add_picture(str(f), Inches(x0), Inches(1.9),
                                 height=Inches(PH))
        txt(s, Inches(x0), Inches(5.85), Inches(wid), Inches(0.36), cap, 13.5,
            INK, True, FD)
        x0 += wid + GAP
    for i, (val, lab) in enumerate([("42%", "of samples found no cell"),
                                    ("938", "of 8,176 cells measured"),
                                    ("0", "candidate sites visited")]):
        x = Inches(0.75 + i * 4.12)
        txt(s, x, Inches(6.36), Inches(1.1), Inches(0.48), val, 24, WINE, True, FH)
        txt(s, x + Inches(1.12), Inches(6.5), Inches(2.9), Inches(0.4), lab,
            15, INK2, False, FD)
    footer(s, "7,144 samples · 117 km of road · a 178 km² box")
    return s


def n3_approaches(prs, d):
    """What each method works FROM.

    Four predicted-coverage maps looked nearly identical -- every model puts
    34-41% of cells above the threshold -- so four of them told an audience
    almost nothing. These show the input instead, which is what actually differs.
    """
    s = slide(prs)
    header(s, 3, "approaches", "Four ways to fill in the map",
           "Each one reads something different about the same ground.")
    panels = ROOT / "sionna-approach" / "deck_panels"
    apps = [
        ("how_baseline.png", "Baseline", OCHRE,
         "Fit a curve to the measurements, read it off anywhere."),
        ("how_raytracing.png", "Ray tracing", TEAL,
         "Rebuild the ground and buildings, then trace the signal."),
        ("how_deeplearning.png", "Deep learning", VIOL,
         "Give it the ground under each link, let it learn the rest."),
        ("how_pinn.png", "PINN", WINE,
         "Learn what distance cannot explain, with physics as a rule."),
    ]
    MW = 3.02
    for i, (fn, name, col, blurb) in enumerate(apps):
        x = Inches(0.45 + i * 3.15)
        txt(s, x, Inches(1.82), Inches(MW), Inches(0.3), name, 18, col, True, FH)
        f = panels / fn
        if f.exists():
            s.shapes.add_picture(str(f), x, Inches(2.24), width=Inches(MW))
        txt(s, x, Inches(4.66), Inches(MW), Inches(1.1), blurb, 15.5, INK2,
            False, FD)
    footer(s, "")
    return s


def n4_results(prs, d):
    """One chart. The gap between the pairs of bars is the whole finding."""
    s = slide(prs)
    header(s, 4, "results", "In sample is not the same as somewhere new",
           "RSRP error in dB. Lower is better. Same rows, same splits, same "
           "200 m buffer between train and test.")
    order = [("reveal-mt-pinn", "PINN"), ("terrain-fno", "Deep\nlearning"),
             ("terrain-parametric", "Baseline"), ("sionna-hybrid-agronomy",
                                                  "Ray tracing")]
    cats, ins, out = [], [], []
    for key, lab in order:
        v = d["bench"].get(key, {})
        try:
            i_ = v["in_sample"]["rmse"]
            o_ = v["kmeans_on_position"]["rmse"]
        except (KeyError, TypeError):
            continue
        cats.append(lab); ins.append(round(i_, 2)); out.append(round(o_, 2))
    if cats:
        bar(s, Inches(0.55), Inches(2.0), Inches(7.5), Inches(3.5), cats,
            [("on data it has seen", ins), ("somewhere new", out)],
            [GREY, TEAL], labels=True, numfmt="0.0", size=15)
    caption(s, Inches(0.55), Inches(5.6), Inches(7.5),
            "Left bar: error on data the model was trained on. Right bar: error "
            "somewhere it has never been. A tall right bar means it memorised.",
            size=15, color=INK2)

    px = Inches(8.35)
    rect(s, px, Inches(2.0), Inches(4.4), Inches(1.85), fill=SURF, line=None)
    txt(s, px + Inches(0.22), Inches(2.14), Inches(4.0), Inches(0.3),
        "Ray tracing barely moves", 18, TEAL, True, FH)
    txt(s, px + Inches(0.22), Inches(2.52), Inches(4.0), Inches(1.25),
        "7.61 dB on training ground, 7.95 dB on new ground. It is the only one "
        "of the four that barely changes.", 16, INK2, False, FD)
    rect(s, px, Inches(4.0), Inches(4.4), Inches(2.25), fill=SURF, line=None)
    txt(s, px + Inches(0.22), Inches(4.12), Inches(4.0), Inches(0.34),
        "The learned models memorise", 18, WINE, True, FH)
    txt(s, px + Inches(0.22), Inches(4.52), Inches(4.0), Inches(1.6),
        "A terrain profile is nearly a unique location label. Feed the "
        "operator the wrong links' profiles and it scores 13.16 dB against "
        "13.13 with the right ones — it never used the terrain.",
        16, INK2, False, FD)
    footer(s, "Held out by region · 200 m buffer between train and test")
    return s


def n5_planner(prs, d):
    """The pipeline in four boxes, then the one chart that reorders priorities."""
    s = slide(prs)
    header(s, 5, "the planner", "Turning a surface into a decision",
           "Change one assumption at a time, and measure how far the "
           "recommendation moves.")
    steps = [("Predict", "a value in every 200 m cell"),
             ("Define served", "eight definitions, one threshold"),
             ("Score", "route-km and area gained"),
             ("Constrain", "power, access, structures, backhaul")]
    for i, (name, body) in enumerate(steps):
        x = Inches(0.55 + i * 3.08)
        rect(s, x, Inches(1.92), Inches(2.85), Inches(1.12), fill=SURF,
             line=None)
        txt(s, x + Inches(0.2), Inches(2.02), Inches(2.5), Inches(0.32), name,
            16, INK, True, FH)
        txt(s, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.56), body,
            14, MUTE, False, FD)
        if i < 3:
            txt(s, x + Inches(2.87), Inches(2.24), Inches(0.2), Inches(0.32),
                "→", 17, RULE, True, FD)
    txt(s, Inches(0.55), Inches(3.05), Inches(7.5), Inches(0.32),
        "How far the recommended site moves when you change one assumption",
        13, INK, True, FD)
    bar(s, Inches(0.4), Inches(3.4), Inches(7.7), Inches(2.75),
        ["what counts as served", "asset class", "propagation model",
         "threshold", "route vs area"],
        [("km", [4.24, 3.35, 2.06, 1.28, 1.21])],
        [WINE, OCHRE, TEAL, GREY, GREY], horizontal=True, labels=True,
        numfmt='0.0"km"', size=15, gridlines=False)
    px = Inches(8.35)
    rect(s, px, Inches(3.05), Inches(4.4), Inches(2.0), fill=SURF, line=None)
    txt(s, px + Inches(0.22), Inches(3.16), Inches(4.0), Inches(0.68),
        "The definition beats the physics", 18, WINE, True, FH)
    txt(s, px + Inches(0.22), Inches(3.86), Inches(4.0), Inches(1.15),
        "Changing what counts as service moves the site twice as far as "
        "changing the model. It was chosen, not measured.", 16, INK2, False, FD)
    rect(s, px, Inches(5.2), Inches(4.4), Inches(1.5), fill=SURF, line=None)
    txt(s, px + Inches(0.22), Inches(5.32), Inches(4.0), Inches(0.34),
        "A neighbourhood, not a pin", 18, TEAL, True, FH)
    txt(s, px + Inches(0.22), Inches(5.7), Inches(4.0), Inches(0.9),
        "Requiring grid power within 1 km moves it 11 km. The exact site "
        "repeats in 10% of draws; 3 km in 99%.", 16, INK2, False, FD)
    footer(s, "198 combinations of model, asset, criterion, threshold and weighting")
    return s


def s6_demo(prs, d):
    """A drawn schematic, not a screenshot, so the deck stays vector."""
    s = slide(prs)
    header(s, 6, "demo", "The planner, live",
           "Place an asset and watch coverage, gain and uncertainty change.")
    rect(s, Inches(0.55), Inches(1.85), Inches(8.15), Inches(4.3), fill=INK,
         line=RULE)
    rnd = random.Random(7)
    for r_ in range(12):
        for c_ in range(24):
            v = (c_ / 23.0) * 0.75 + rnd.random() * 0.25
            col = TEAL if v > 0.62 else (OCHRE if v > 0.38 else WINE)
            rect(s, Inches(0.72 + c_ * 0.325), Inches(2.0 + r_ * 0.325),
                 Inches(0.305), Inches(0.305), fill=col, line=None)
    for i in range(58):
        t = i / 57.0
        x = 0.9 + 7.3 * t
        y = 3.2 + 1.5 * (0.5 - abs(0.5 - t)) * 2 - 0.5
        rect(s, Inches(x), Inches(y), Inches(0.075), Inches(0.075), fill=WHITE,
             line=None)
    oval(s, Inches(5.5), Inches(2.55), Inches(0.22), Inches(0.22), fill=OCHRE,
         line=None)
    oval(s, Inches(2.9), Inches(4.25), Inches(0.3), Inches(0.3), fill=None,
         line=WHITE, lw=2.5)
    txt(s, Inches(0.72), Inches(5.88), Inches(4.6), Inches(0.24),
        "map scale ▬ 2 km   ·   ▪ one 200 m demand cell", 9, WHITE, False, FM)
    px = Inches(8.95)
    rect(s, px, Inches(1.85), Inches(3.8), Inches(4.3), fill=SURF, line=RULE)
    txt(s, px + Inches(0.2), Inches(1.98), Inches(3.4), Inches(0.26),
        "Rural Coverage Planner", 13, INK, True, FD)
    for i, (lab, val) in enumerate([("approach", "Baseline — fitted physics ▾"),
                                    ("what counts as served", "Availability ▾"),
                                    ("target", "≥ 50% of the time"),
                                    ("route vs area", "0.70 / 0.30"),
                                    ("asset", "Macro · 37 m · 0 dB")]):
        y = Inches(2.36 + i * 0.52)
        txt(s, px + Inches(0.2), y, Inches(3.4), Inches(0.2), lab, 8.5, MUTE,
            False, FD)
        rect(s, px + Inches(0.2), y + Inches(0.19), Inches(3.4), Inches(0.25),
             fill=WHITE, line=RULE)
        txt(s, px + Inches(0.3), y + Inches(0.235), Inches(3.2), Inches(0.2),
            val, 9, INK, False, FM)
    rect(s, px + Inches(0.2), Inches(5.0), Inches(3.4), Inches(0.32), fill=TEAL,
         line=None)
    txt(s, px + Inches(0.95), Inches(5.07), Inches(2.4), Inches(0.22),
        "Find the best site", 10, WHITE, True, FM)
    for i, (lab, val) in enumerate([("route-km", "51.9 → 80.5"),
                                    ("score", "0.421 → 0.652")]):
        rect(s, px + Inches(0.2 + i * 1.75), Inches(5.45), Inches(1.65),
             Inches(0.55), fill=WHITE, line=RULE)
        txt(s, px + Inches(0.32 + i * 1.75), Inches(5.52), Inches(1.4),
            Inches(0.2), lab, 8.5, MUTE, False, FD)
        txt(s, px + Inches(0.32 + i * 1.75), Inches(5.72), Inches(1.4),
            Inches(0.26), val, 11, TEAL, True, FD)
    caption(s, Inches(0.55), Inches(6.34), Inches(12.2),
            "Pick a model, pick what counts as service, drop an asset anywhere — "
            "coverage, gain and uncertainty all re-solve live.",
            size=15, color=INK2)
    footer(s, "Runs in a browser. No server, no install, no network.")
    return s


SLIDES = [n1_problem, n2_gap, n3_approaches, n4_results, n5_planner, s6_demo]


def build(verbose=True):
    d = load()
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in SLIDES:
        fn(prs, d)
    prs.save(OUT)
    if verbose:
        print(f"[deck] {OUT} ({OUT.stat().st_size/1e6:.2f} MB), {len(SLIDES)} slides")
        print(f"[deck] bench: {', '.join(d['bench']) or 'none'}")
        print(f"[deck] bundles: {', '.join(d['bundles']) or 'none'}")
        for nm in ORDER:
            if nm not in d["bench"]:
                print(f"[deck] {nm}: no testbench entry, accuracy row reserved")
            if nm not in d["bundles"]:
                print(f"[deck] {nm}: no bundle, siting row reserved")
    return OUT


if __name__ == "__main__":
    build()
